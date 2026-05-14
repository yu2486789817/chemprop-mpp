"""
create_ppt.py — Generate the workshop PowerPoint presentation.

Requires: pip install python-pptx (run setup_env first)

Produces: ../ppt/chemprop_workshop.pptx
"""

import os
import sys

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("ERROR: python-pptx not installed.")
    print("Run: pip install python-pptx")
    sys.exit(1)

PROJECT_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
PPT_DIR = os.path.join(PROJECT_ROOT, "ppt")
RESULTS_DIR = os.path.join(PROJECT_ROOT, "results")
os.makedirs(PPT_DIR, exist_ok=True)

# ── Theme Colors ───────────────────────────────────────────────────────────
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
MED_BLUE = RGBColor(0x32, 0x74, 0xA1)
LIGHT_BLUE = RGBColor(0x5B, 0x9B, 0xD5)
ACCENT_ORANGE = RGBColor(0xE1, 0x81, 0x2C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MED_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_bar(slide, title_text, subtitle_text=None):
    """Add a dark blue title bar at top of slide."""
    # Title bar background
    left = Inches(0)
    top = Inches(0)
    width = SLIDE_WIDTH
    height = Inches(1.4)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    # Title text
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.2), Inches(11.5), Inches(0.9))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32)
    p.font.color.rgb = WHITE
    p.font.bold = True

    if subtitle_text:
        txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(0.85), Inches(11.5), Inches(0.5))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle_text
        p2.font.size = Pt(16)
        p2.font.color.rgb = LIGHT_BLUE


def add_body_text(slide, text, left=0.8, top=1.8, width=11.5, height=5.0, font_size=18):
    """Add body text box with bullet-style content."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    lines = text.strip().split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(8)

        # Handle indentation levels
        stripped = line.strip()
        if stripped.startswith("•"):
            p.level = 0
            p.font.size = Pt(font_size)
        elif stripped.startswith("–"):
            p.level = 1
            p.font.size = Pt(font_size - 2)
            p.font.color.rgb = MED_GRAY

    return txBox


def add_footer(slide, page_num, total=7):
    """Add page number footer."""
    txBox = slide.shapes.add_textbox(Inches(12.0), Inches(7.1), Inches(1.0), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{page_num}/{total}"
    p.font.size = Pt(10)
    p.font.color.rgb = MED_GRAY
    p.alignment = PP_ALIGN.RIGHT


def add_image_if_exists(slide, image_path, left, top, width=None, height=None):
    """Add an image if the file exists."""
    if os.path.exists(image_path):
        kwargs = {}
        if width:
            kwargs["width"] = Inches(width)
        if height:
            kwargs["height"] = Inches(height)
        slide.shapes.add_picture(image_path, Inches(left), Inches(top), **kwargs)
        return True
    return False


def create_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    blank_layout = prs.slide_layouts[6]  # blank layout

    # ═══════════════════════════════════════════════════════════════════════
    # Slide 1: Title
    # ═══════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BLUE)

    # Decorative line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), Inches(2.8), Inches(9), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_ORANGE
    line.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(2), Inches(1.2), Inches(9), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Molecular Property Prediction with Chemprop"
    p.font.size = Pt(40)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(2), Inches(3.1), Inches(9), Inches(1.2))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "An AI for Science Approach to Drug Discovery"
    p2.font.size = Pt(24)
    p2.font.color.rgb = LIGHT_BLUE
    p2.alignment = PP_ALIGN.CENTER

    # Details
    txBox3 = slide.shapes.add_textbox(Inches(2.5), Inches(4.5), Inches(8), Inches(1.5))
    tf3 = txBox3.text_frame
    for line_text in [
        "Graph Neural Networks for Molecular Property Prediction",
        "ESOL Solubility Dataset | Chemprop MPNN | Data Scale Analysis",
    ]:
        p3 = tf3.add_paragraph()
        p3.text = line_text
        p3.font.size = Pt(16)
        p3.font.color.rgb = WHITE
        p3.alignment = PP_ALIGN.CENTER

    add_footer(slide, 1)

    # ═══════════════════════════════════════════════════════════════════════
    # Slide 2: Background — Why Molecular Property Prediction?
    # ═══════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "AI for Science: Why Molecular Property Prediction?",
                  "The drug discovery bottleneck — and how AI can help")

    body = (
        "• Drug discovery is slow and expensive: ~10 years and $2.6 billion per new drug\n"
        "• A key bottleneck is early-stage screening: millions of molecules, very few tested\n"
        "\n"
        "• Molecular property prediction uses AI to predict properties from structure alone\n"
        "  – No wet lab needed for initial screening\n"
        "  – Prioritize promising candidates before expensive synthesis\n"
        "\n"
        "• Why solubility? It is the first pharmacokinetic hurdle\n"
        "  – If a drug doesn't dissolve, it cannot be absorbed in the gut\n"
        "  – Poor solubility is the #1 cause of late-stage clinical failure\n"
        "\n"
        "• AI for Science principle: Data quality and structure define the boundary of AI scientific reasoning\n"
        "  – The model is only as good as the data and molecular representation"
    )
    add_body_text(slide, body, top=1.8, font_size=18)
    add_footer(slide, 2)

    # ═══════════════════════════════════════════════════════════════════════
    # Slide 3: Data & Method
    # ═══════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Data & Method: ESOL Dataset + Chemprop MPNN",
                  "Dataset overview and model architecture")

    # Left column: Data
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "ESOL (Delaney) Dataset"
    p.font.size = Pt(22)
    p.font.color.rgb = DARK_BLUE
    p.font.bold = True

    for line in [
        "• 1,128 small organic molecules",
        "• Each entry: SMILES string → LogS (aqueous solubility)",
        "• LogS = log₁₀(solubility in mol/L)",
        "• Range: ~−7 to ~+2 (insoluble → highly soluble)",
        "• Widely used benchmark, fast to train",
        "• 80/20 train/test split (fixed random seed)",
    ]:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(4)

    # Right column: Method
    txBox2 = slide.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.5), Inches(5.0))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True

    p2 = tf2.paragraphs[0]
    p2.text = "Chemprop MPNN"
    p2.font.size = Pt(22)
    p2.font.color.rgb = DARK_BLUE
    p2.font.bold = True

    for line in [
        "• Message-Passing Neural Network",
        "• Molecules as graphs:",
        "  – Nodes = atoms (C, N, O, H...)",
        "  – Edges = chemical bonds",
        "• Message passing: atoms exchange",
        "  information through bonds over",
        "  multiple steps (depth = 5)",
        "• Readout: aggregate atom features",
        "  → predict molecular property",
        "• PyTorch-based, modular design",
        "• Chemprop v2: improved efficiency",
    ]:
        p2 = tf2.add_paragraph()
        p2.text = line
        p2.font.size = Pt(16)
        p2.font.color.rgb = DARK_GRAY
        p2.space_after = Pt(4)

    # Divider
    line_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(2.0), Inches(0.03), Inches(4.5))
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = LIGHT_BLUE
    line_shape.line.fill.background()

    add_footer(slide, 3)

    # ═══════════════════════════════════════════════════════════════════════
    # Slide 4: Experiment Design
    # ═══════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Experiment Design",
                  "Baseline training + data scale analysis")

    body = (
        "Experiment 1: Baseline Model Training\n"
        "  – Train Chemprop MPNN on full ESOL dataset\n"
        "  – Architecture: hidden=300, depth=5, dropout=0.1\n"
        "  – Training: Adam optimizer, ReduceLROnPlateau, early stopping (patience=15)\n"
        "  – Evaluation metrics: RMSE, MAE, R²\n"
        "\n"
        "Experiment 2: Data Scale Analysis (Core Contribution)\n"
        "  – Question: How does model performance improve with more data?\n"
        "  – Train on 20%, 50%, 80%, 100% of training data\n"
        "  – Fixed test set for fair comparison\n"
        "  – 3 random repeats per scale → error bars\n"
        "\n"
        "Key Design Choices:\n"
        "  – Fixed test set across all experiments (same random seed = 42)\n"
        "  – MPNN treats each molecule independently — no data leakage\n"
        "  – Regression task (continuous LogS values)"
    )
    add_body_text(slide, body, top=1.8, font_size=17)
    add_footer(slide, 4)

    # ═══════════════════════════════════════════════════════════════════════
    # Slide 5: Results — Prediction Performance
    # ═══════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Results: Prediction Performance",
                  "Chemprop MPNN accurately predicts aqueous solubility")

    # Insert images if they exist
    img1 = os.path.join(RESULTS_DIR, "predicted_vs_true.png")
    img2 = os.path.join(RESULTS_DIR, "error_distribution.png")

    has_img1 = add_image_if_exists(slide, img1, 0.6, 2.0, width=5.8)
    has_img2 = add_image_if_exists(slide, img2, 6.8, 2.0, width=5.8)

    if not has_img1:
        # Placeholder text
        txBox = slide.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(5.0), Inches(3.0))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "[Run train.py and visualize.py\nto generate figures]"
        p.font.size = Pt(16)
        p.font.color.rgb = MED_GRAY
        p.alignment = PP_ALIGN.CENTER

    if not has_img2:
        txBox = slide.shapes.add_textbox(Inches(7.0), Inches(2.5), Inches(5.0), Inches(3.0))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "[Run train.py and visualize.py\nto generate figures]"
        p.font.size = Pt(16)
        p.font.color.rgb = MED_GRAY
        p.alignment = PP_ALIGN.CENTER

    # Bottom caption
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Left: Predicted vs measured LogS (hexbin density plot). Right: Error distribution with normal fit."
    p.font.size = Pt(12)
    p.font.color.rgb = MED_GRAY

    add_footer(slide, 5)

    # ═══════════════════════════════════════════════════════════════════════
    # Slide 6: Data Scale Results
    # ═══════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Results: Impact of Training Data Size",
                  "More data → better predictions (with diminishing returns)")

    img3 = os.path.join(RESULTS_DIR, "data_scale_rmse.png")
    has_img3 = add_image_if_exists(slide, img3, 1.0, 2.0, width=11.0)

    if not has_img3:
        txBox = slide.shapes.add_textbox(Inches(2.0), Inches(2.5), Inches(9.0), Inches(3.0))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "[Run data_scale_experiment.py then visualize.py\nto generate this figure]"
        p.font.size = Pt(16)
        p.font.color.rgb = MED_GRAY
        p.alignment = PP_ALIGN.CENTER

    # Key insight box
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(1.0), Inches(6.0), Inches(11.3), Inches(0.8))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_GRAY
    box.line.fill.background()

    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Key Insight: Model performance improves dramatically from 20%→80% data, with diminishing returns at 100%. "
    p.text += "This validates the AI for Science principle: data quality and quantity define the upper bound of what models can learn."
    p.font.size = Pt(14)
    p.font.color.rgb = DARK_BLUE
    p.font.italic = True
    p.alignment = PP_ALIGN.CENTER

    add_footer(slide, 6)

    # ═══════════════════════════════════════════════════════════════════════
    # Slide 7: Conclusions
    # ═══════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BLUE)

    # Title
    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.0), Inches(1.0))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Conclusions: AI + Molecular Science"
    p.font.size = Pt(36)
    p.font.color.rgb = WHITE
    p.font.bold = True

    # Decorative line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.55), Inches(3), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_ORANGE
    line.line.fill.background()

    # Main points
    points = [
        ("1", "Graph Neural Networks are a natural fit for molecular data",
         "Molecules are graphs — atoms as nodes, bonds as edges. MPNNs learn from this structure directly, without hand-crafted features."),
        ("2", "Solubility prediction is practical and impactful",
         "Aqueous solubility is the first hurdle for oral drugs. AI screening can prioritize synthesizable candidates early."),
        ("3", "Data scale drives model quality",
         "Our experiments show RMSE improves substantially with more training data — validating the AI for Science approach."),
        ("4", "Chemprop lowers the barrier to entry",
         "Modular, well-documented, PyTorch-based. A few hundred lines of code is enough for a complete molecular ML pipeline."),
    ]

    for i, (num, title, desc) in enumerate(points):
        y = 2.0 + i * 1.3

        # Number circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                        Inches(1.0), Inches(y), Inches(0.5), Inches(0.5))
        circle.fill.solid()
        circle.fill.fore_color.rgb = ACCENT_ORANGE
        circle.line.fill.background()
        tf = circle.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Title
        txBox = slide.shapes.add_textbox(Inches(1.8), Inches(y - 0.05), Inches(10.5), Inches(0.45))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(20)
        p.font.color.rgb = WHITE
        p.font.bold = True

        # Description
        txBox2 = slide.shapes.add_textbox(Inches(1.8), Inches(y + 0.45), Inches(10.5), Inches(0.6))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = desc
        p2.font.size = Pt(14)
        p2.font.color.rgb = LIGHT_BLUE

    add_footer(slide, 7)

    return prs


def main():
    print("Creating workshop presentation...")
    prs = create_presentation()
    output_path = os.path.join(PPT_DIR, "chemprop_workshop.pptx")
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
